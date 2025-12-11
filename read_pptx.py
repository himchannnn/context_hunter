from pptx import Presentation
import os

filename = "설계 문서.pptx"

if not os.path.exists(filename):
    print(f"File not found: {filename}")
    exit(1)

try:
    prs = Presentation(filename)
    with open("pptx_content.txt", "w", encoding="utf-8") as f:
        f.write(f"--- Extracting text from {filename} ---\n")
        
        for i, slide in enumerate(prs.slides):
            f.write(f"\n[Slide {i+1}]\n")
            text_content = []
            if slide.shapes.title:
                f.write(f"Title: {slide.shapes.title.text}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                    text_content.append(shape.text)
            
            if text_content:
                f.write("\n".join(text_content) + "\n")
    print("Done writing to pptx_content.txt")
except Exception as e:
    print(f"Error reading pptx: {e}")
